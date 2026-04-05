from .base import BaseParser, ParsedDocument
import io


class PPTXParser(BaseParser):
    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text = ""

            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n=== Slide {slide_num} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            return ParsedDocument(
                texto=text,
                estrutura=self._detect_structure(text),
                metadados={
                    "arquivo": filename,
                    "tamanho": len(content),
                    "slides": len(prs.slides),
                },
                alertas=[],
                limitacoes=["Imagens e gráficos podem não ser extraídos"],
            )
        except Exception as e:
            result = self._create_empty_result(filename)
            result.alertas.append(f"Erro ao processar PPTX: {str(e)}")
            return result


class TxtParser(BaseParser):
    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            text = content.decode("utf-8")
            return ParsedDocument(
                texto=text,
                estrutura=self._detect_structure(text),
                metadados={
                    "arquivo": filename,
                    "tamanho": len(content),
                    "encoding": "utf-8",
                },
                alertas=[],
                limitacoes=[],
            )
        except Exception as e:
            try:
                text = content.decode("latin-1")
                return ParsedDocument(
                    texto=text,
                    estrutura=self._detect_structure(text),
                    metadados={
                        "arquivo": filename,
                        "tamanho": len(content),
                        "encoding": "latin-1",
                    },
                    alertas=["Arquivo codificado em latin-1"],
                    limitacoes=[],
                )
            except:
                result = self._create_empty_result(filename)
                result.alertas.append(f"Erro ao processar TXT: {str(e)}")
                return result


class MdParser(BaseParser):
    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            text = content.decode("utf-8")
            return ParsedDocument(
                texto=text,
                estrutura=self._detect_structure(text),
                metadados={
                    "arquivo": filename,
                    "tamanho": len(content),
                    "formato": "markdown",
                },
                alertas=[],
                limitacoes=[],
            )
        except Exception as e:
            result = self._create_empty_result(filename)
            result.alertas.append(f"Erro ao processar MD: {str(e)}")
            return result
